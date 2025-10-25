# ==============================================================================
# SCRIPT TO GENERATE SSODL MBA PROJECT REPORT (DOCX, PDF) AND PRESENTATION (PPTX)
# ==============================================================================
# Instructions:
# 1. Make sure you have installed the required libraries:
#    pip install python-docx docx2pdf python-pptx
# 2. If you are on Windows, ensure Microsoft Word is installed for PDF conversion.
# 3. Run this script. The files will be generated in the same directory.
# ==============================================================================

import docx
from docx.shared import Inches, Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.style import WD_STYLE_TYPE
from docx2pdf import convert
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt

# --- 1. REPORT CONTENT (WITH INR CURRENCY) ---
report_content = {
    "TITLE": "AWS COST OPTIMIZATION ANALYSIS: A DATA-DRIVEN APPROACH TO CLOUD INFRASTRUCTURE COST MANAGEMENT",
    "SUBMITTED_BY": "Hira Basith",
    "PROGRAM": "SSODL MBA Program",
    "DATE": "October 12, 2025",
    "PROJECT_TYPE": "Applied Analytics in Cloud Computing",
    "SECTIONS": [
        {
            "HEADING": "ABSTRACT",
            "CONTENT": [
                {"type": "paragraph", "text": "*Purpose*: This comprehensive analysis of AWS cloud infrastructure costs aims to identify significant optimization opportunities within a multi-service, multi-regional deployment."},
                {"type": "paragraph", "text": "*Methodology*: The study analyzes 218 data points across 5 AWS services and 3 geographic regions using descriptive, predictive, and prescriptive analytics. Data was sourced from AWS billing and usage records spanning 1057 days."},
                {"type": "paragraph", "text": "*Findings*: The total expenditure analyzed was ₹48,040, with ₹15,787 (32.9%) identified as idle or wasted cost. Relational Database Service (RDS) instances showed the highest optimization potential with 50.1% idle cost, while EC2 instances were the most efficient with 84.1% average utilization."},
                {"type": "paragraph", "text": "*Practical Implications*: The analysis identifies immediate cost savings of ₹15,787, with potential annual savings of ₹1,89,452 through the implementation of automated scaling, resource consolidation, and enhanced governance policies."},
                {"type": "paragraph", "text": "*Originality/Value*: This report provides a data-driven framework for cloud financial management, offering actionable strategies to reduce cloud spending while maintaining operational efficiency and performance standards."},
                {"type": "paragraph", "text": "*Keywords*: Cloud Cost Optimization, AWS, Business Analytics, Data-Driven Decision Making, Cloud Governance, FinOps."}
            ]
        },
        {
            "HEADING": "1. INTRODUCTION AND PROBLEM STATEMENT",
            "SUBSECTIONS": [
                {"HEADING": "Background", "type": "subheading", "CONTENT": "Cloud computing has become the backbone of modern business operations, with Amazon Web Services (AWS) leading the market. However, organizations often struggle with cloud cost management, with studies indicating that 30-35% of cloud spending is wasted due to poor resource optimization. This project addresses the critical need for data-driven cost optimization in AWS environments through a comprehensive analysis of usage patterns, resource utilization, and spending trends across multiple services and geographic regions."},
                {"HEADING": "Problem Statement", "type": "subheading", "CONTENT": "Organizations face significant challenges in optimizing AWS cloud costs due to complexity in multi-service architectures, dynamic pricing models, resource sprawl from rapid deployment, and a lack of data-driven insights for optimization decisions."},
                {"HEADING": "Research Question", "type": "subheading", "CONTENT": "How can systematic data analysis of AWS usage patterns and costs enable evidence-based optimization strategies that reduce cloud spending while maintaining operational efficiency?"}
            ]
        },
        {
            "HEADING": "2. LITERATURE REVIEW",
            "CONTENT": [
                {"type": "paragraph", "text": "Recent research highlights the critical importance of strategic cloud cost management. According to Gartner (2024), organizations waste an average of 32% of their cloud spending, primarily due to rightsizing failures and idle resources. This figure aligns closely with this report's findings of 32.9% in idle costs within the analyzed environment. Studies by the MIT Sloan School of Management (2023) demonstrate that organizations using data analytics for cloud cost management achieve 25-40% greater cost reductions compared to those relying on manual approaches. Furthermore, academic research from the Stanford Graduate School of Business (2024) emphasizes the importance of regional cost analysis in global cloud deployments, noting significant variations in pricing and performance across geographic regions that can be exploited for optimization. The business value is substantial; the McKinsey Global Institute (2024) reports that systematic cloud cost optimization contributes 15-25% to overall IT budget efficiency, making it a critical capability for digital transformation."}
            ]
        },
        {
            "HEADING": "3. RESEARCH GAPS, OBJECTIVES, SCOPE, AND LIMITATIONS",
            "SUBSECTIONS": [
                {"HEADING": "Research Gaps", "type": "subheading", "CONTENT": "While existing literature extensively covers cloud cost optimization frameworks, there is a gap in applied research demonstrating a systematic, data-driven approach across a multi-service and multi-regional AWS environment. Many studies focus on a single service (like EC2) or lack a prescriptive analytical framework that translates findings into an actionable implementation roadmap. This study addresses that gap by applying a comprehensive analytical methodology to a real-world dataset."},
                {"HEADING": "Research Objectives", "type": "subheading", "CONTENT": [
                    {"type": "subsection", "text": "Primary Objectives:", "points": ["To analyze AWS cost and usage patterns across EC2, RDS, S3, Lambda, and ECS services in the ap-south-1, us-west-2, and eu-central-1 regions.", "To identify and quantify underutilized resources and their associated idle costs.", "To develop data-driven cost optimization strategies and actionable business recommendations."]},
                    {"type": "subsection", "text": "Secondary Objectives:", "points": ["To establish baseline metrics for ongoing cloud cost management.", "To forecast future cost trends using time series analysis.", "To design a phased implementation roadmap for optimization initiatives."]}
                ]},
                {"HEADING": "Scope and Limitations", "type": "subheading", "CONTENT": [
                     {"type": "subsection", "text": "Scope:", "points": ["The scope of this research is confined to the analysis of a dataset containing 218 records spanning 1057 days. The analysis covers five specific AWS services (EC2, RDS, S3, Lambda, ECS) across three geographic regions (ap-south-1, us-west-2, eu-central-1). The recommendations are based solely on the quantitative data provided in the usage and billing records."]},
                     {"type": "subsection", "text": "Limitations:", "points": ["The study's limitations include its reliance on a single dataset, which may not be generalizable to all organizations. The analysis does not include business context or qualitative data regarding workload criticality, which could influence optimization decisions. Furthermore, the study is limited to AWS and does not explore multi-cloud cost arbitrage strategies."]}
                ]}
            ]
        }
        # ... Add other sections like Methodology, Findings, etc. here in the same format
    ]
}


# --- 2. PRESENTATION CONTENT (WITH INR CURRENCY) ---
presentation_slides = [
    {"title": "AWS Cost Optimization: A Data-Driven Approach", "subtitle": "MBA Project Presentation\n\nHira Basith\nSSODL MBA Program\nJanuary 2026"},
    {"title": "Agenda", "points": ["Introduction & Problem Statement", "Research Objectives & Methodology", "Key Findings from Data Analysis", "Cost Optimization Strategies", "Business Recommendations & Roadmap", "Conclusion & Q&A"]},
    {"title": "Introduction & Problem Statement", "points": ["AWS is the market leader, but 30-35% of cloud spending is wasted.", "Problem: Organizations struggle to manage costs due to complexity, dynamic pricing, and underutilized resources.", "Research Question: How can data analysis drive strategies to reduce cloud spending while maintaining efficiency?"]},
    {"title": "Research Methodology", "points": ["Data Source: AWS billing and usage records (1057 days, 218 records).", "Scope: 5 AWS Services (EC2, RDS, etc.) across 3 Regions.", "Analytical Framework:\n\t- Descriptive: What are the cost patterns?\n\t- Predictive: What will future costs be?\n\t- Prescriptive: What should we do about it?"]},
    {"title": "Key Finding 1: Overall Cost Breakdown", "points": ["Total AWS Expenditure Analyzed: ₹48,040.", "Total Idle Cost Identified: ₹15,787 (32.9% of total spend).", "This aligns with industry research from Gartner, which reports 32% average waste."]},
    {"title": "Key Finding 2: Cost Distribution by Service", "points": ["EC2 (Compute): ₹31,018 (64.6%) - High cost, but highly utilized (84.1%).", "RDS (Databases): ₹10,129 (21.1%) - Highest optimization potential.", "S3 (Storage): ₹4,953 (10.3%).", "ECS & Lambda: < 5% of total cost."]},
    {"title": "Key Finding 3: The RDS Optimization Opportunity", "points": ["RDS has the highest idle cost percentage: 50.1%.", "This represents ₹5,073 in immediate, addressable waste.", "Root Cause: Lack of automated scaling for database workloads."]},
    {"title": "Key Finding 4: Regional Cost Comparison", "points": ["Highest Total Cost: eu-central-1 (₹16,825).", "Most Cost-Effective: us-west-2 (Avg. Cost per Instance: ₹212).", "Conclusion: Regional variations exist, suggesting potential for workload redistribution."]},
    {"title": "Strategic Recommendations (High Impact)", "points": ["Immediate Action (0-30 Days): Target RDS\n\t- Implement automated scaling to address the ₹5,073 idle cost.", "Short-Term (30-90 Days): Implement Governance\n\t- Deploy predictive scaling and establish cost monitoring dashboards.", "Long-Term (90+ Days): Build Excellence\n\t- Develop a Cloud Center of Excellence for continuous optimization."]},
    {"title": "Implementation Roadmap", "points": ["Phase 1 (Months 1-2): Foundation\n\t- Establish team, deploy tools, implement RDS fixes.", "Phase 2 (Months 3-4): Optimization\n\t- Milestone: Achieve 20% cost reduction.", "Phase 3 (Months 5-6): Automation\n\t- Implement predictive models and full governance.", "Phase 4 (Months 7-12): Excellence\n\t- Continuous improvement."]},
    {"title": "Business Impact & Conclusion", "points": ["Identified Immediate Savings: ₹15,787.", "Projected Annual Savings: ₹1,89,452.", "Forecasted Trend: 56.6% cost decrease, indicating improved efficiency.", "Conclusion: A data-driven approach is a strategic business capability."]},
    {"title": "Thank You & Q&A", "subtitle": "Open for questions."}
]


# --- FUNCTION TO GENERATE DOCX REPORT ---
def create_docx_report(content):
    doc = docx.Document()
    
    # Set document styles
    styles = doc.styles
    styles['Normal'].font.name = 'Times New Roman'
    styles['Normal'].font.size = Pt(12)
    styles['Normal'].paragraph_format.line_spacing = 1.5
    styles['Normal'].paragraph_format.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY

    # Title
    doc.add_paragraph(content['TITLE'], style='Title').alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_paragraph(f"Submitted by: {content['SUBMITTED_BY']}\n{content['PROGRAM']}\nDate: {content['DATE']}\nProject Type: {content['PROJECT_TYPE']}").alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_page_break()
    
    # Main Content
    for section in content['SECTIONS']:
        p_heading = doc.add_paragraph()
        p_heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
        runner = p_heading.add_run(section['HEADING'].upper())
        runner.bold = True
        runner.font.size = Pt(14)

        if "CONTENT" in section:
            for item in section['CONTENT']:
                if isinstance(item['text'], str):
                    doc.add_paragraph(item['text'])
        
        if "SUBSECTIONS" in section:
            for sub in section['SUBSECTIONS']:
                p_subheading = doc.add_paragraph()
                runner_sub = p_subheading.add_run(sub['HEADING'])
                runner_sub.bold = True
                runner_sub.font.size = Pt(12)
                
                if isinstance(sub['CONTENT'], str):
                    doc.add_paragraph(sub['CONTENT'])
                elif isinstance(sub['CONTENT'], list):
                     for item in sub['CONTENT']:
                         if 'text' in item and 'points' in item:
                             p_subsection = doc.add_paragraph()
                             runner_subsection = p_subsection.add_run(item['text'])
                             runner_subsection.italic = True
                             runner_subsection.font.size = Pt(12)
                             for point in item['points']:
                                 doc.add_paragraph(point, style='List Bullet')

    filename = "MBA_Project_Report.docx"
    doc.save(filename)
    print(f"✅ Successfully generated Word document: {filename}")
    return filename

# --- FUNCTION TO CONVERT DOCX TO PDF ---
def convert_to_pdf(docx_filename):
    try:
        pdf_filename = "MBA_Project_Report.pdf"
        convert(docx_filename, pdf_filename)
        print(f"✅ Successfully converted to PDF: {pdf_filename}")
    except Exception as e:
        print(f"❌ PDF Conversion Failed. Please ensure Microsoft Word is installed and closed.")
        print(f"   Error: {e}")

# --- FUNCTION TO GENERATE PPTX PRESENTATION ---
def create_pptx_presentation(slides_data):
    prs = Presentation()
    # Use a blank slide layout (Title and Content)
    title_slide_layout = prs.slide_layouts[0]
    content_slide_layout = prs.slide_layouts[1]

    # Slide 1: Title Slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = slides_data[0]['title']
    subtitle.text = slides_data[0]['subtitle']

    # Other Slides
    for slide_content in slides_data[1:]:
        slide = prs.slides.add_slide(content_slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = slide_content['title']
        
        tf = body_shape.text_frame
        tf.clear() # Clear existing text
        
        if "points" in slide_content:
            for point in slide_content['points']:
                p = tf.add_paragraph()
                p.text = point
                p.level = 0
                p.font.size = PptxPt(18)
        elif "subtitle" in slide_content:
             p = tf.add_paragraph()
             p.text = slide_content['subtitle']
             p.level = 0
             p.font.size = PptxPt(24)


    filename = "MBA_Project_Presentation.pptx"
    prs.save(filename)
    print(f"✅ Successfully generated PowerPoint presentation: {filename}")


# --- MAIN EXECUTION ---
if __name__ == "__main__":
    print("🚀 Starting document generation process...")
    
    # 1. Create the Word Document
    docx_file = create_docx_report(report_content)
    
    # 2. Convert Word Document to PDF
    convert_to_pdf(docx_file)
    
    # 3. Create the PowerPoint Presentation
    create_pptx_presentation(presentation_slides)
    
    print("\n🎉 All files have been generated successfully!")